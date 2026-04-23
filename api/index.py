import os
import io
import base64
import json
import urllib.request
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import anthropic

app = Flask(__name__, static_folder=os.path.join(os.path.dirname(__file__), '..', 'public'), static_url_path='')
CORS(app)

UNSPLASH_ACCESS_KEY = os.environ.get('UNSPLASH_ACCESS_KEY', '')


def hex2rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def get_colours(style):
    s = (style or '').lower()
    if 'bright' in s or 'colour' in s or 'color' in s:
        return dict(dk='#1A1A2E', h1='#7C4DFF', h2='#E040FB', ac='#FF6D00', tx='#1A1A2E', lt='#FFF8E1', wh='#FFFFFF', text_on_lt='#1A1A2E')
    if 'minimal' in s or 'white' in s or 'clean' in s:
        return dict(dk='#0D1B2A', h1='#1565C0', h2='#0D47A1', ac='#00ACC1', tx='#1A1A2E', lt='#E8F4FD', wh='#FFFFFF', text_on_lt='#1A1A2E')
    if 'pastel' in s or 'friendly' in s:
        return dict(dk='#5B4B6B', h1='#C2467A', h2='#9B59B6', ac='#1ABC9C', tx='#2C3E50', lt='#FDEEF4', wh='#FFFFFF', text_on_lt='#2C3E50')
    if 'projector' in s or 'contrast' in s:
        return dict(dk='#000000', h1='#000000', h2='#222222', ac='#FFD600', tx='#000000', lt='#FFFFFF', wh='#FFFFFF', text_on_lt='#000000')
    return dict(dk='#1A1D2E', h1='#1E2761', h2='#6C63FF', ac='#22C55E', tx='#1E2030', lt='#EEF0FF', wh='#FFFFFF', text_on_lt='#1E2030')


def fetch_image(query):
    """Fetch an image from Unsplash based on query. Returns image bytes or None."""
    try:
        if not UNSPLASH_ACCESS_KEY:
            return None
        url = f"https://api.unsplash.com/photos/random?query={urllib.parse.quote(query)}&orientation=landscape&client_id={UNSPLASH_ACCESS_KEY}"
        req = urllib.request.Request(url, headers={'User-Agent': 'Preperly/1.0'})
        with urllib.request.urlopen(req, timeout=5) as resp:
            data = json.loads(resp.read())
            img_url = data['urls']['regular']
        req2 = urllib.request.Request(img_url, headers={'User-Agent': 'Preperly/1.0'})
        with urllib.request.urlopen(req2, timeout=8) as img_resp:
            return img_resp.read()
    except Exception:
        return None


def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex2rgb(color)
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=16, color='#FFFFFF', bold=False, italic=False, align='left', font='Century Gothic'):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(text or '')
    run.font.size = Pt(size)
    run.font.color.rgb = hex2rgb(color)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb


def add_bullets(slide, items, x, y, w, h, size=14, color='#1E2030', font='Century Gothic'):
    if not items:
        return
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = '• ' + str(item)
        run.font.size = Pt(size)
        run.font.color.rgb = hex2rgb(color)
        run.font.name = font
    return tb


def add_image_to_slide(slide, img_bytes, x, y, w, h):
    """Add image bytes to slide at given position."""
    try:
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, Inches(x), Inches(y), Inches(w), Inches(h))
        return True
    except Exception:
        return False


def build_pptx(plan, slides):
    C = get_colours(plan.get('slideStyle', ''))
    font = 'Century Gothic'

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    def new_slide(bg):
        s = prs.slides.add_slide(blank)
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = hex2rgb(bg)
        return s

    objs  = (plan.get('learningObjectives') or [])[:4]
    vocab = (plan.get('keyVocabulary') or [])[:4]
    hook  = slides.get('hook') or {}
    c1    = slides.get('content1') or {}
    c2    = slides.get('content2') or {}
    act   = slides.get('activity') or {}
    cp    = slides.get('continuousProvision') or {}

    subject = plan.get('subject', '')
    topic   = plan.get('lessonTitle', '')

    # Try to fetch a topic image
    topic_img = fetch_image(f"{subject} {topic} children education")

    # ── Slide 1: Title ──
    s = new_slide(C['dk'])
    # Decorative accent bar on left
    add_rect(s, 0, 0, 0.25, 5.625, C['h2'])
    add_rect(s, 0.25, 0, 9.75, 0.08, C['ac'])
    # If we have an image, put it on the right half with overlay
    if topic_img:
        add_rect(s, 5, 0, 5, 5.625, C['h1'])
        add_image_to_slide(s, topic_img, 5, 0, 5, 5.625)
        # Dark overlay on image
        overlay = add_rect(s, 5, 0, 5, 5.625, C['dk'])
        overlay.fill.fore_color.rgb = hex2rgb(C['dk'])
        try:
            from pptx.util import Pt as PPt
            from pptx.oxml.ns import qn
            from lxml import etree
            sp_elem = overlay._element
            sp_pr = sp_elem.find(qn('p:spPr'))
            solidFill = sp_pr.find(qn('a:solidFill'))
            srgbClr = solidFill.find(qn('a:srgbClr'))
            alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha_elem.set('val', '75000')
        except Exception:
            pass
        txt_w = 4.2
    else:
        txt_w = 9.0

    add_text(s, subject.upper(), 0.5, 0.5, txt_w, 0.5, size=11, color=C['h2'], bold=True, font=font)
    add_text(s, topic, 0.5, 1.1, txt_w, 2.2, size=28, color='#FFFFFF', bold=True, font=font)
    add_rect(s, 0.5, 3.4, 1.8, 0.07, C['ac'])
    add_text(s, f"{plan.get('yearGroup','')}  •  {plan.get('duration','')}", 0.5, 3.6, txt_w, 0.5, size=14, color='#B0B8D0', font=font)

    # ── Slide 2: Objectives ──
    s = new_slide(C['lt'])
    add_rect(s, 0, 0, 10, 1.0, C['h1'])
    add_rect(s, 0, 1.0, 10, 0.06, C['ac'])
    add_text(s, '🎯  LEARNING OBJECTIVES', 0.4, 0.2, 9, 0.65, size=18, color='#FFFFFF', bold=True, font=font)
    add_text(s, 'By the end of this lesson, you will be able to:', 0.4, 1.2, 9.2, 0.45, size=13, color=C['tx'], italic=True, font=font)
    for i, obj in enumerate(objs):
        y = 1.8 + i * 0.88
        # Numbered circle background
        add_rect(s, 0.4, y + 0.05, 0.6, 0.6, C['h2'])
        add_text(s, str(i + 1), 0.4, y + 0.05, 0.6, 0.6, size=16, color='#FFFFFF', bold=True, align='center', font=font)
        add_rect(s, 1.15, y, 8.4, 0.72, C['wh'])
        add_text(s, obj, 1.3, y + 0.08, 8.1, 0.58, size=13, color=C['tx'], font=font)

    # ── Slide 3: Starter / Hook ──
    s = new_slide(C['dk'])
    add_rect(s, 0, 0, 10, 0.08, C['ac'])
    add_rect(s, 0, 5.545, 10, 0.08, C['ac'])
    add_text(s, '⚡  STARTER ACTIVITY', 0.5, 0.2, 9, 0.5, size=12, color=C['ac'], bold=True, font=font)
    add_text(s, hook.get('title') or 'Think About This...', 0.5, 0.75, 9, 1.0, size=26, color='#FFFFFF', bold=True, font=font)
    add_rect(s, 0.5, 1.85, 9, 0.05, C['h2'])
    bullets = hook.get('bullets') or []
    for i, b in enumerate(bullets[:4]):
        add_rect(s, 0.5, 2.05 + i * 0.85, 0.5, 0.6, C['h2'])
        add_text(s, '★', 0.5, 2.1 + i * 0.85, 0.5, 0.5, size=14, color='#FFFFFF', align='center', font=font)
        add_text(s, b, 1.15, 2.08 + i * 0.85, 8.4, 0.65, size=14, color='#E8EAF2', font=font)

    # ── Slide 4: Content 1 ──
    s = new_slide(C['wh'])
    add_rect(s, 0, 0, 10, 1.0, C['h1'])
    add_rect(s, 0, 1.0, 0.08, 4.625, C['ac'])
    add_text(s, c1.get('title') or '', 0.35, 0.15, 9.3, 0.75, size=20, color='#FFFFFF', bold=True, font=font)
    bullets = c1.get('bullets') or []
    for i, b in enumerate(bullets[:5]):
        add_rect(s, 0.25, 1.15 + i * 0.82, 0.35, 0.35, C['h2'])
        add_text(s, '→', 0.25, 1.18 + i * 0.82, 0.35, 0.32, size=11, color='#FFFFFF', align='center', font=font)
        add_text(s, b, 0.75, 1.15 + i * 0.82, 9.0, 0.7, size=14, color=C['tx'], font=font)

    # ── Slide 5: Content 2 ──
    s = new_slide(C['lt'])
    add_rect(s, 0, 0, 10, 1.0, C['h2'])
    add_rect(s, 0, 1.0, 0.08, 4.625, C['h1'])
    add_text(s, c2.get('title') or '', 0.35, 0.15, 9.3, 0.75, size=20, color='#FFFFFF', bold=True, font=font)
    bullets = c2.get('bullets') or []
    for i, b in enumerate(bullets[:5]):
        add_rect(s, 0.25, 1.15 + i * 0.82, 0.35, 0.35, C['h1'])
        add_text(s, '→', 0.25, 1.18 + i * 0.82, 0.35, 0.32, size=11, color='#FFFFFF', align='center', font=font)
        add_text(s, b, 0.75, 1.15 + i * 0.82, 9.0, 0.7, size=14, color=C['tx'], font=font)

    # ── Slide 6: Vocabulary ──
    s = new_slide(C['lt'])
    add_rect(s, 0, 0, 10, 1.0, C['h2'])
    add_text(s, '📖  KEY VOCABULARY', 0.4, 0.18, 9, 0.65, size=18, color='#FFFFFF', bold=True, font=font)
    for i, v in enumerate(vocab):
        col = i % 2
        row = i // 2
        x = 0.3 if col == 0 else 5.2
        y = 1.1 + row * 2.1
        # Card
        add_rect(s, x, y, 4.6, 1.85, C['wh'])
        add_rect(s, x, y, 4.6, 0.55, C['h1'])
        add_text(s, (v.get('word') or '').upper(), x + 0.15, y + 0.08, 4.3, 0.42, size=14, color='#FFFFFF', bold=True, font=font)
        add_text(s, v.get('definition') or '', x + 0.15, y + 0.65, 4.3, 1.1, size=12, color=C['tx'], font=font)

    # ── Slide 7: Activity ──
    s = new_slide(C['dk'])
    add_rect(s, 0, 0, 10, 0.08, C['h2'])
    add_text(s, '✏️  ACTIVITY TIME', 0.5, 0.18, 9, 0.5, size=12, color=C['ac'], bold=True, font=font)
    add_text(s, act.get('title') or 'Your Task', 0.5, 0.72, 9, 0.9, size=24, color='#FFFFFF', bold=True, font=font)
    if act.get('time'):
        add_rect(s, 0.5, 1.7, 2.8, 0.52, C['ac'])
        add_text(s, '⏱  ' + act['time'], 0.5, 1.75, 2.8, 0.44, size=14, color='#FFFFFF', bold=True, align='center', font=font)
    steps = (act.get('steps') or [])[:4]
    for i, step in enumerate(steps):
        add_rect(s, 0.5, 2.38 + i * 0.77, 0.55, 0.55, C['h2'])
        add_text(s, str(i+1), 0.5, 2.42 + i * 0.77, 0.55, 0.48, size=15, color='#FFFFFF', bold=True, align='center', font=font)
        add_text(s, step, 1.2, 2.4 + i * 0.77, 8.5, 0.65, size=14, color='#E8EAF2', font=font)

    # ── Slide 8: Summary ──
    s = new_slide(C['wh'])
    add_rect(s, 0, 0, 10, 1.0, C['h1'])
    add_text(s, '✅  LESSON SUMMARY', 0.4, 0.18, 9, 0.65, size=18, color='#FFFFFF', bold=True, font=font)
    add_text(s, 'What have we learned today?', 0.4, 1.1, 9, 0.42, size=13, color='#8890B0', italic=True, font=font)
    for i, obj in enumerate(objs):
        y = 1.65 + i * 0.88
        add_rect(s, 0.4, y, 9.2, 0.72, C['lt'])
        add_rect(s, 0.4, y, 0.55, 0.72, C['ac'])
        add_text(s, '✓', 0.4, y + 0.12, 0.55, 0.48, size=14, color='#FFFFFF', bold=True, align='center', font=font)
        add_text(s, obj, 1.1, y + 0.1, 8.3, 0.55, size=13, color=C['tx'], font=font)

    # ── Slide 9: Exit Ticket ──
    s = new_slide(C['dk'])
    add_rect(s, 0, 0, 10, 0.08, C['h2'])
    add_text(s, '🎫  EXIT TICKET', 0.5, 0.2, 9, 0.5, size=12, color=C['ac'], bold=True, font=font)
    add_text(s, 'Before You Go...', 0.5, 0.78, 9, 0.8, size=24, color='#FFFFFF', bold=True, font=font)
    # Question box
    add_rect(s, 0.5, 1.72, 9, 2.0, C['h2'])
    add_rect(s, 0.5, 1.72, 9, 0.08, C['ac'])
    add_text(s, slides.get('exitQ') or "Summarise today's lesson in 3 key points.", 0.7, 1.9, 8.6, 1.65, size=16, color='#FFFFFF', italic=True, font=font)
    # Answer lines
    for i in range(3):
        add_rect(s, 0.5, 3.9 + i * 0.5, 9, 0.04, '#3A4060')
    add_text(s, f"{subject}  •  {plan.get('yearGroup','')}  •  {topic}", 0.5, 5.22, 9, 0.3, size=8, color='#4A5070', font=font)

    # ── Slide 10: Continuous Provision (EYFS/Year 1 only) ──
    year = (plan.get('yearGroup') or '').lower()
    if cp and ('eyfs' in year or 'reception' in year or 'year 1' in year or 'year1' in year):
        s = new_slide('#1A3A2A')
        add_rect(s, 0, 0, 10, 1.0, '#2D6A4F')
        add_text(s, '🌱  CONTINUOUS PROVISION', 0.4, 0.18, 9, 0.65, size=18, color='#FFFFFF', bold=True, font=font)
        areas = cp.get('areas') or []
        for i, area in enumerate(areas[:4]):
            col = i % 2
            row = i // 2
            x = 0.3 if col == 0 else 5.2
            y = 1.15 + row * 2.05
            add_rect(s, x, y, 4.55, 1.85, '#FFFFFF')
            add_rect(s, x, y, 4.55, 0.52, '#52B788')
            add_text(s, area.get('name') or '', x + 0.15, y + 0.08, 4.25, 0.38, size=13, color='#FFFFFF', bold=True, font=font)
            add_text(s, area.get('activity') or '', x + 0.15, y + 0.62, 4.25, 1.15, size=11, color='#1A3A2A', font=font)
        if cp.get('adultFocus'):
            add_rect(s, 0.3, 5.08, 9.4, 0.42, '#52B788')
            add_text(s, '👩‍🏫 Adult Focus: ' + cp['adultFocus'], 0.5, 5.1, 9.0, 0.38, size=11, color='#FFFFFF', bold=True, font=font)

    buf = io.BytesIO()
    prs.save(buf)
    return base64.b64encode(buf.getvalue()).decode('utf-8')


# ── Routes ────────────────────────────────────────────────────

@app.route('/')
def index():
    return send_from_directory(app.static_folder, 'index.html')

@app.route('/<path:path>')
def static_files(path):
    return send_from_directory(app.static_folder, path)

@app.route('/api/generate-plan', methods=['POST'])
def generate_plan():
    try:
        data         = request.json
        api_key      = data.get('apiKey', '')
        if not api_key:
            return jsonify({'error': 'No API key'}), 400

        subject      = data.get('subject', '')
        year_group   = data.get('yearGroup', '')
        topic        = data.get('topic', '')
        duration     = data.get('duration', '60 minutes')
        requirements = data.get('requirements', 'none')

        # Detect if EYFS/Year 1 for continuous provision
        yg = year_group.lower()
        is_early_years = 'eyfs' in yg or 'reception' in yg or 'year 1' in yg

        cp_instruction = ''
        if is_early_years:
            cp_instruction = ',"continuousProvision":{"areas":[{"name":"Area name e.g. Maths Area","activity":"What children do here linked to lesson topic"},{"name":"...","activity":"..."},{"name":"...","activity":"..."},{"name":"...","activity":"..."}],"adultFocus":"What the adult-led activity will be"}'

        client = anthropic.Anthropic(api_key=api_key)
        message = client.messages.create(
            model='claude-sonnet-4-6',
            max_tokens=4000,
            messages=[{
                'role': 'user',
                'content': (
                    'RESPOND WITH ONLY VALID JSON. NO BACKTICKS. NO EXPLANATION. START WITH { END WITH }.\n'
                    f'Create a UK school lesson plan: Subject:{subject}, Year:{year_group}, '
                    f'Topic:{topic}, Duration:{duration}, Notes:{requirements or "none"}\n'
                    'JSON:{"lessonTitle":"...","learningObjectives":["...","...","..."],'
                    '"keyVocabulary":[{"word":"...","definition":"..."},{"word":"...","definition":"..."},{"word":"...","definition":"..."}],'
                    '"resources":["...","..."],'
                    '"lessonPhases":[{"phase":"Starter","duration":"10 mins","teacherActivity":"...","studentActivity":"...","differentiation":"..."},'
                    '{"phase":"Direct Teaching","duration":"15 mins","teacherActivity":"...","studentActivity":"...","differentiation":"..."},'
                    '{"phase":"Main Activity","duration":"20 mins","teacherActivity":"...","studentActivity":"...","differentiation":"..."},'
                    '{"phase":"Plenary","duration":"5 mins","teacherActivity":"...","studentActivity":"...","differentiation":"..."}],'
                    '"assessmentStrategies":["...","..."],"homeworkSuggestion":"...","teacherNotes":"..."' +
                    cp_instruction + '}'
                )
            }]
        )
        raw = message.content[0].text
        # Strip markdown code blocks if present
        raw = raw.replace('```json', '').replace('```', '').strip()
        s = raw.index('{')
        e = raw.rindex('}')
        plan = json.loads(raw[s:e+1])
        plan.update({'subject': subject, 'yearGroup': year_group, 'duration': duration})
        return jsonify({'success': True, 'plan': plan})

    except Exception as ex:
        import traceback
        return jsonify({'success': False, 'error': str(ex), 'trace': traceback.format_exc()}), 500


@app.route('/api/build-pptx', methods=['POST'])
def build_pptx_route():
    try:
        data    = request.json
        api_key = data.get('apiKey', '')
        plan    = data.get('plan', {})
        if not api_key:
            return jsonify({'error': 'No API key'}), 400

        client = anthropic.Anthropic(api_key=api_key)
        objs   = (plan.get('learningObjectives') or [])[:3]

        yg = (plan.get('yearGroup') or '').lower()
        is_early_years = 'eyfs' in yg or 'reception' in yg or 'year 1' in yg

        cp_instruction = ''
        if is_early_years:
            cp_instruction = ',"continuousProvision":{"areas":[{"name":"...","activity":"..."},{"name":"...","activity":"..."},{"name":"...","activity":"..."},{"name":"...","activity":"..."}],"adultFocus":"..."}'

        message = client.messages.create(
            model='claude-sonnet-4-6',
            max_tokens=2000,
            messages=[{
                'role': 'user',
                'content': (
                    'RESPOND WITH ONLY VALID JSON. NO BACKTICKS. NO EXPLANATION. START WITH { END WITH }.\n'
                    f'Slide text for: "{plan.get("lessonTitle","")}" ({plan.get("subject","")}, {plan.get("yearGroup","")}).\n'
                    f'Objectives: {"; ".join(objs)}\n'
                    'JSON:{"hook":{"title":"...","bullets":["...","...","..."]},'
                    '"content1":{"title":"...","bullets":["...","...","..."]},'
                    '"content2":{"title":"...","bullets":["...","...","..."]},'
                    '"activity":{"title":"Activity: ...","time":"20 mins","steps":["...","...","..."]},'
                    '"exitQ":"..."' + cp_instruction + '}'
                )
            }]
        )
        raw    = message.content[0].text
        # Strip markdown code blocks if present
        raw = raw.replace('```json', '').replace('```', '').strip()
        s      = raw.index('{')
        e      = raw.rindex('}')
        slides = json.loads(raw[s:e+1])

        # Pass continuous provision from plan if present
        if plan.get('continuousProvision') and not slides.get('continuousProvision'):
            slides['continuousProvision'] = plan['continuousProvision']

        pptx_b64 = build_pptx(plan, slides)
        return jsonify({'success': True, 'pptx': pptx_b64, 'slides': slides})

    except Exception as ex:
        return jsonify({'success': False, 'error': str(ex)}), 500


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port)
